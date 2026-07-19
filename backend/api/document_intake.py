"""Safe, bounded extraction for decks and other founder-provided materials.

Only uploaded files are read.  We retain extracted text and per-page/slide metadata,
not raw scraped pages.  Visual previews are generated from the submitted file and
served locally so the investor can inspect the exact source behind a claim.
"""
from __future__ import annotations

import re
import shutil
from pathlib import Path
from typing import Any


SUPPORTED_EXTENSIONS = {".pdf", ".pptx", ".docx", ".txt", ".md"}


class DocumentIntakeError(ValueError):
    pass


def allowed_filename(filename: str) -> bool:
    return Path(filename).suffix.lower() in SUPPORTED_EXTENSIONS


def extract_document(path: Path, media_root: Path, company_id: str) -> tuple[str, list[dict[str, Any]]]:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return _extract_pdf(path, media_root, company_id)
    if suffix == ".pptx":
        return _extract_pptx(path, media_root, company_id)
    if suffix == ".docx":
        return _extract_docx(path, media_root, company_id)
    if suffix in {".txt", ".md"}:
        text = path.read_text(encoding="utf-8", errors="replace")
        return text, [{"kind": "document", "title": path.name, "page": 1, "text": text[:4000]}]
    raise DocumentIntakeError("Supported files are PDF, PPTX, DOCX, TXT, and Markdown.")


def _extract_pdf(path: Path, media_root: Path, company_id: str) -> tuple[str, list[dict[str, Any]]]:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise DocumentIntakeError("PDF extraction requires pypdf. Run pip install -r requirements.txt.") from exc

    reader = PdfReader(str(path))
    pages: list[dict[str, Any]] = []
    chunks: list[str] = []
    preview_dir = media_root / "previews" / company_id
    for number, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        chunks.append(f"=== SLIDE {number} ===\n{text}")
        pages.append({"kind": "deck_slide", "title": f"Deck slide {number}", "page": number, "text": text[:6000]})

    # Render previews when PyMuPDF is available. Extraction still works without it.
    try:
        import fitz

        preview_dir.mkdir(parents=True, exist_ok=True)
        pdf = fitz.open(str(path))
        for number, page in enumerate(pdf, start=1):
            preview = preview_dir / f"slide-{number}.png"
            page.get_pixmap(matrix=fitz.Matrix(1.25, 1.25), alpha=False).save(str(preview))
            pages[number - 1]["preview_url"] = f"/media/previews/{company_id}/{preview.name}"
    except ImportError:
        pass
    return "\n\n".join(chunks), pages


def _extract_pptx(path: Path, media_root: Path, company_id: str) -> tuple[str, list[dict[str, Any]]]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise DocumentIntakeError("PPTX extraction requires python-pptx. Run pip install -r requirements.txt.") from exc

    presentation = Presentation(str(path))
    slides: list[dict[str, Any]] = []
    chunks: list[str] = []
    image_dir = media_root / "extracted-images" / company_id
    for number, slide in enumerate(presentation.slides, start=1):
        slide_text = []
        image_urls = []
        for index, shape in enumerate(slide.shapes, start=1):
            if getattr(shape, "has_text_frame", False):
                slide_text.append(shape.text)
            if getattr(shape, "shape_type", None) == 13:  # MSO_SHAPE_TYPE.PICTURE
                image_dir.mkdir(parents=True, exist_ok=True)
                extension = shape.image.ext or "png"
                target = image_dir / f"slide-{number}-image-{index}.{extension}"
                target.write_bytes(shape.image.blob)
                image_urls.append(f"/media/extracted-images/{company_id}/{target.name}")
        text = "\n".join(part for part in slide_text if part.strip())
        chunks.append(f"=== SLIDE {number} ===\n{text}")
        slides.append({"kind": "deck_slide", "title": f"Deck slide {number}", "page": number, "text": text[:6000], "image_urls": image_urls})
    return "\n\n".join(chunks), slides


def _extract_docx(path: Path, media_root: Path, company_id: str) -> tuple[str, list[dict[str, Any]]]:
    try:
        from docx import Document
    except ImportError as exc:
        raise DocumentIntakeError("DOCX extraction requires python-docx. Run pip install -r requirements.txt.") from exc
    document = Document(str(path))
    text = "\n".join(paragraph.text for paragraph in document.paragraphs if paragraph.text.strip())
    image_dir = media_root / "extracted-images" / company_id
    image_urls = []
    for index, relationship in enumerate(document.part.rels.values(), start=1):
        if "image" not in relationship.reltype:
            continue
        image_dir.mkdir(parents=True, exist_ok=True)
        extension = Path(relationship.target_ref).suffix or ".png"
        target = image_dir / f"image-{index}{extension}"
        target.write_bytes(relationship.target_part.blob)
        image_urls.append(f"/media/extracted-images/{company_id}/{target.name}")
    return text, [{"kind": "document", "title": path.name, "page": 1, "text": text[:12000], "image_urls": image_urls}]


def copy_uploaded_file(source, destination: Path, max_bytes: int) -> int:
    """Copy an UploadFile stream while enforcing the body limit."""
    total = 0
    with destination.open("wb") as output:
        while chunk := source.file.read(1024 * 1024):
            total += len(chunk)
            if total > max_bytes:
                output.close()
                destination.unlink(missing_ok=True)
                raise DocumentIntakeError("File is larger than the 25 MB upload limit.")
            output.write(chunk)
    return total


def copy_photo(source, destination: Path, max_bytes: int) -> int:
    content_type = (source.content_type or "").lower()
    if not content_type.startswith("image/"):
        raise DocumentIntakeError("Founder photo must be an image file.")
    return copy_uploaded_file(source, destination, max_bytes)
